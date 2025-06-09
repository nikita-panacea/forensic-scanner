# import pdfplumber
# from docx import Document

# def extract_text_from_pdf(file_path):
#     texts = []
#     with pdfplumber.open(file_path) as pdf:
#         for i, page in enumerate(pdf.pages):
#             texts.append((i + 1, page.extract_text()))
#     return texts

# def extract_text_from_docx(file_path):
#     doc = Document(file_path)
#     full_text = "\n".join(p.text for p in doc.paragraphs)
#     return [(1, full_text)]

import fitz
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import pandas as pd
from ebooklib import epub
from odf import text, teletype
from odf.opendocument import load as load_odf
import extract_msg
import pytesseract
from PIL import Image
import zipfile
import tempfile
import os

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    for i, page in enumerate(doc, 1):
        yield f"PDF page {i}", page.get_text()

def extract_text_from_docx(path):
    doc = Document(path)
    # paragraphs
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    # tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    yield "DOCX", "\n".join(parts)

def extract_text_from_pptx(path):
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        yield f"PPTX slide {i}", "\n".join(texts)

def extract_text_from_html(path):
    html = open(path, encoding="utf-8", errors="ignore").read()
    soup = BeautifulSoup(html, "lxml")
    yield "HTML", soup.get_text(separator="\n")

def extract_text_from_xlsx(path):
    df = pd.read_excel(path, engine="openpyxl", header=None, dtype=str)
    yield "XLSX", "\n".join(df.fillna("").astype(str).values.flatten())

def extract_text_from_csv(path):
    df = pd.read_csv(path, header=None, dtype=str)
    yield "CSV", "\n".join(df.fillna("").astype(str).values.flatten())

def extract_text_from_txt(path):
    yield "TXT", open(path, encoding="utf-8", errors="ignore").read()

def extract_text_from_epub(path):
    book = epub.read_epub(path)
    parts = []
    for item in book.get_items():
        if item.get_type() == epub.EpubHtml:
            soup = BeautifulSoup(item.get_content(), "lxml")
            parts.append(soup.get_text(separator="\n"))
    yield "EPUB", "\n".join(parts)

def extract_text_from_odt(path):
    odt = load_odf(path)
    text_parts = [teletype.extractText(el) for el in odt.getElementsByType(text.P)]
    yield "ODT", "\n".join(text_parts)

def extract_text_from_msg(path):
    msg = extract_msg.Message(path)
    if msg.subject:
        yield "MSG subject", msg.subject
    if msg.body:
        yield "MSG body", msg.body

def extract_text_from_image(path):
    img = Image.open(path)
    text = pytesseract.image_to_string(img)
    yield "OCR image", text

def extract_text_from_zip(path):
    """
    Extract all entries from the zip (or jar, odt, etc.),
    then feed each file inside back through the main extractor map.
    """
    from config import EXTRACTORS  # avoid circular import

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            with zipfile.ZipFile(path, "r") as z:
                for member in z.namelist():
                    # Skip directories
                    if member.endswith("/"):
                        continue
                    # Extract to temp
                    extracted_path = z.extract(member, tmpdir)
                    # Determine extension
                    ext = os.path.splitext(member)[1].lower()
                    extractor = EXTRACTORS.get(ext)
                    if extractor:
                        try:
                            for label, text in extractor(extracted_path):
                                yield f"ZIP:{member}:{label}", text
                        except Exception:
                            # Skip entries that fail
                            continue
        except zipfile.BadZipFile:
            # Not a valid zip
            return